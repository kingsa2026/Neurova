"""
附件文本抽取器（R-3）

把上传附件转换为可注入 LLM 上下文的文本：
- text（txt/md/rst/json/yaml/toml/log 等文本族）/ code → 直接按 UTF-8/gb18030 解码
- document：docx（python-docx 段落+表格）、xlsx（openpyxl 单元格）、pptx（python-pptx 文本+备注）、
  pdf（pypdf 页面文字）、csv → 表格文本化、html/htm/xml → 剥标签、rtf（striprtf）、
  odt/ods/odp（odfpy）
- image / audio / video → 不抽取文本（由调用方走 vision/ASR/占位）
- 其他（含 ppt/doc/xls 旧格式）→ 返回 None（调用方降级为元数据占位）

所有解析失败都降级为 None（返回 (None, reason) 供日志排查），不抛异常——附件处理失败
不拖垮整轮对话。文件大小不设限；真正防超大文档拖垮上下文的闸门是
MAX_EXTRACT_CHARS（抽取文本上限）。
"""

import io
import os
from typing import Optional, Tuple

MAX_EXTRACT_CHARS = 200_000  # 20 万字符上限（真正的上下文闸门）


def _decode_text(data: bytes, filename: str) -> Optional[str]:
    text = None
    for enc in ("utf-8", "gb18030", "latin-1"):
        try:
            text = data.decode(enc)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    return text[:MAX_EXTRACT_CHARS] if text else None


def _extract_docx(data: bytes) -> str:
    import docx  # lazy import：依赖缺失不阻断其他类型

    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)[:MAX_EXTRACT_CHARS]


def _extract_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[工作表: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
            if sum(len(p) for p in parts) > MAX_EXTRACT_CHARS:
                break
    wb.close()
    return "\n".join(parts)[:MAX_EXTRACT_CHARS]


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"[幻灯片 {idx}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        # 备注也可以抽取
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[备注] {notes}")
        except Exception:
            pass
    return "\n".join(parts)[:MAX_EXTRACT_CHARS]


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages[:100]:  # 最多 100 页，防超大 PDF
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""
        if page_text.strip():
            parts.append(page_text)
    return "\n".join(parts)[:MAX_EXTRACT_CHARS]


def extract_attachment_text(data: bytes, filename: str, file_type: str) -> Tuple[Optional[str], str]:
    """按扩展名抽取附件文本。

    Returns:
        (text, status) —— text 为抽取结果（不可用时 None），status 为诊断原因。
    """
    if data is None or len(data) == 0:
        return None, "empty_file"

    ext = os.path.splitext(filename)[1].lower()

    if file_type in ("text", "code"):
        text = _decode_text(data, filename)
        return (text, "decoded") if text else (None, "decode_failed")

    if ext == ".docx":
        try:
            return _extract_docx(data), "docx"
        except Exception as e:
            return None, f"docx_error:{type(e).__name__}"
    if ext == ".xlsx":
        try:
            return _extract_xlsx(data), "xlsx"
        except Exception as e:
            return None, f"xlsx_error:{type(e).__name__}"
    if ext == ".pptx":
        try:
            return _extract_pptx(data), "pptx"
        except Exception as e:
            return None, f"pptx_error:{type(e).__name__}"
    if ext == ".pdf":
        try:
            return _extract_pdf(data), "pdf"
        except Exception as e:
            return None, f"pdf_error:{type(e).__name__}"
    if ext == ".csv":
        text = _decode_text(data, filename)
        return (text, "csv") if text else (None, "csv_decode_failed")
    if ext in (".html", ".htm", ".xml"):
        try:
            return _extract_html(data), "html"
        except Exception as e:
            return None, f"html_error:{type(e).__name__}"
    if ext == ".rtf":
        try:
            return _extract_rtf(data), "rtf"
        except Exception as e:
            return None, f"rtf_error:{type(e).__name__}"
    if ext in (".odt", ".ods", ".odp"):
        try:
            return _extract_odf(data), ext.lstrip(".")
        except Exception as e:
            return None, f"odf_error:{type(e).__name__}"

    # 未支持的文档/二进制（ppt/doc/xls 旧格式）不抽取文本。
    # 注意 .ppt（OLE2）此前被错误归入 pptx 分支，python-pptx 打不开只会抛
    # pptx_error 静默失败；python-pptx 仅支持 OOXML（.pptx）。
    return None, "unsupported_format"


def _extract_rtf(data: bytes) -> str:
    """RTF → 纯文本（striprtf；依赖缺失回退 stdlib 正则剥控制字）。"""
    try:
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(data.decode("ascii", errors="ignore"))[:MAX_EXTRACT_CHARS]
    except ImportError:
        import re

        text = data.decode("ascii", errors="ignore")
        text = re.sub(r"\\par[d]?\b", "\n", text)
        text = re.sub(r"\{\\\*?[^{}]*\}", "", text)  # 剥 {\*\...} 目标组
        text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)  # 十六进制转义剥除
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
        return re.sub(r"[{}]", "", text)[:MAX_EXTRACT_CHARS]


def _extract_odf(data: bytes) -> str:
    """ODF 家族（odt 文本 / ods 表格 / odp 演示）→ 纯文本（odfpy）。"""
    from odf.opendocument import load
    from odf import text as odf_text, table as odf_table, teletype

    doc = load(io.BytesIO(data))
    parts = []
    # 表格单元格优先（ods/odp 里也常有）
    for cell in doc.getElementsByType(odf_table.TableCell):
        cell_text = teletype.extractText(cell)
        if cell_text.strip():
            parts.append(cell_text)
    # 段落与标题
    for elem in doc.getElementsByType(odf_text.P) + doc.getElementsByType(odf_text.H):
        para_text = teletype.extractText(elem)
        if para_text.strip() and para_text not in parts:
            parts.append(para_text)
    return "\n".join(parts)[:MAX_EXTRACT_CHARS]


def _extract_html(data: bytes) -> str:
    """提取 HTML 正文文本（stdlib HTMLParser：剥离 script/style/标签）。"""
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self._skip_depth = 0
            self._parts: list = []

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip_depth += 1

        def handle_endtag(self, tag):
            if tag in ("script", "style") and self._skip_depth > 0:
                self._skip_depth -= 1

        def handle_data(self, data):
            if self._skip_depth == 0:
                text = data.strip()
                if text:
                    self._parts.append(text)

        def get_text(self) -> str:
            return "\n".join(self._parts)

    raw = _decode_text(data, "page.html") or ""
    parser = _TextExtractor()
    try:
        parser.feed(raw)
        parser.close()
    except Exception:
        pass
    text = parser.get_text()
    return text[:MAX_EXTRACT_CHARS] if text else raw[:MAX_EXTRACT_CHARS]
