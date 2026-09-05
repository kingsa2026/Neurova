"""
测试：knowledge API 导入（R-4 多格式文档 + 远程网页）

契约:
  1. POST /knowledge/import 支持文件上传（txt/md/docx/xlsx/pptx/pdf/html），
     抽取文本后创建知识条目，返回条目列表
  2. POST /knowledge/import-url 支持远程网页 URL（http/https），
     抓取 HTML 抽取正文存为条目
  3. SSRF 防护：ftp/file/localhost/环回/私有 IP 一律拒绝
  4. GET /knowledge 无数据时返回空列表（不再返回模拟数据）
"""

import io
import pytest
from unittest.mock import patch

from neurova.api.endpoints import knowledge as kb


def _make_docx_bytes():
    import docx

    d = docx.Document()
    d.add_paragraph("文档正文内容 ABC")
    d.add_paragraph("第二段")
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


class TestImportFile:
    @pytest.mark.asyncio
    async def test_import_txt_creates_item(self, monkeypatch, tmp_path):
        # 注入独立仓库（临时目录）
        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)

        class FakeUpload:
            filename = "note.txt"
            content_type = "text/plain"

            async def read(self):
                return "hello knowledge content".encode("utf-8")

        result = await kb.import_knowledge_file(
            FakeUpload(), request=None, current_user={"user_id": "u1"}, agent_id="default"
        )
        assert result["code"] == 0
        items = result["data"]["items"]
        assert len(items) == 1
        assert "hello knowledge content" in items[0]["content"]

    @pytest.mark.asyncio
    async def test_import_docx_extracts(self, monkeypatch, tmp_path):
        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)

        class FakeUpload:
            filename = "report.docx"
            content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

            async def read(self):
                return _make_docx_bytes()

        result = await kb.import_knowledge_file(
            FakeUpload(), request=None, current_user={"user_id": "u1"}, agent_id="default"
        )
        items = result["data"]["items"]
        assert len(items) == 1
        assert "文档正文内容 ABC" in items[0]["content"]

    @pytest.mark.asyncio
    async def test_import_unsupported_returns_extract_failed(self, monkeypatch, tmp_path):
        """2026-09-06 契约升级：抽取失败不再以成功语义静默返回空列表，
        返回 code=1 + status，前端据此提示用户（修复"提示成功但列表没有"）。"""
        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)

        class FakeUpload:
            filename = "archive.zip"
            content_type = "application/zip"

            async def read(self):
                return b"PK\x05\x06\x00\x00"

        result = await kb.import_knowledge_file(
            FakeUpload(), request=None, current_user={"user_id": "u1"}, agent_id="default"
        )
        assert result["code"] == 1
        assert result["data"]["items"] == []
        assert result["data"]["status"] == "unsupported_format"

    @pytest.mark.asyncio
    async def test_import_legacy_ppt_rejected_as_unsupported(self, monkeypatch, tmp_path):
        """旧版 .ppt（OLE2 二进制）python-pptx 打不开——此前假支持走 pptx 分支
        抛 pptx_error 静默失败；现明确归为不支持格式并显式报错。"""
        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)

        class FakeUpload:
            filename = "legacy.ppt"
            content_type = "application/vnd.ms-powerpoint"

            async def read(self):
                return b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" * 4  # OLE2 魔数

        result = await kb.import_knowledge_file(
            FakeUpload(), request=None, current_user={"user_id": "u1"}, agent_id="default"
        )
        assert result["code"] == 1
        assert result["data"]["status"] == "unsupported_format"
        assert result["data"]["items"] == []

    @pytest.mark.asyncio
    async def test_import_pptx_extracts_slides(self, monkeypatch, tmp_path):
        """真 .pptx 走 python-pptx 抽取：正文+备注，防回归锁定。"""
        from pptx import Presentation

        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "标题甲"
        slide.placeholders[1].text = "要点一"
        buf = io.BytesIO()
        prs.save(buf)

        class FakeUpload:
            filename = "deck.pptx"
            content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

            async def read(self):
                return buf.getvalue()

        result = await kb.import_knowledge_file(
            FakeUpload(), request=None, current_user={"user_id": "u1"}, agent_id="default"
        )
        assert result["code"] == 0
        items = result["data"]["items"]
        assert len(items) == 1
        assert "[幻灯片 1]" in items[0]["content"]
        assert "要点一" in items[0]["content"]

    @pytest.mark.asyncio
    async def test_import_large_file_accepted(self, monkeypatch, tmp_path):
        """2026-09-06 契约升级：2MB 字节限制移除（用户要求），大文件正常入库。"""
        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)

        class FakeUpload:
            filename = "big.md"
            content_type = "text/markdown"

            async def read(self):
                return ("# 标题\n正文段落。\n" * 90000).encode("utf-8")  # >2MB

        result = await kb.import_knowledge_file(
            FakeUpload(), request=None, current_user={"user_id": "u1"}, agent_id="default"
        )
        assert result["code"] == 0
        assert len(result["data"]["items"]) == 1

    @pytest.mark.asyncio
    async def test_import_rtf_and_json(self, monkeypatch, tmp_path):
        """新格式：rtf（striprtf）、json（文本族）。"""
        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)

        for filename, payload in [
            ("doc.rtf", b"{\\rtf1\\ansi RTF body text here.}"),
            ("cfg.json", b'{"key": "json value here"}'),
        ]:
            class FakeUpload:
                content_type = "application/octet-stream"

                async def read(self):
                    return payload

            FakeUpload.filename = filename
            result = await kb.import_knowledge_file(
                FakeUpload(), request=None, current_user={"user_id": "u1"}, agent_id="default"
            )
            assert result["code"] == 0, filename
            assert len(result["data"]["items"]) == 1, filename


class TestImportUrl:
    def test_validate_url_allows_https(self, monkeypatch):
        # 沙箱无外网 DNS——mock getaddrinfo 返回公网 IP，验证公网域名放行
        import socket

        monkeypatch.setattr(
            socket,
            "getaddrinfo",
            lambda host, port: [(socket.AF_INET, socket.SOCK_STREAM, 6, "", ("93.184.216.34", 443))],
        )
        assert kb._validate_import_url("https://example.com/page") is True
        assert kb._validate_import_url("http://example.com") is True

    def test_validate_url_rejects_insecure(self):
        assert kb._validate_import_url("ftp://example.com") is False
        assert kb._validate_import_url("file:///etc/passwd") is False
        assert kb._validate_import_url("javascript:alert(1)") is False
        assert kb._validate_import_url("data:text/html,hi") is False

    def test_validate_url_rejects_localhost_and_private(self):
        assert kb._validate_import_url("http://localhost:8000/x") is False
        assert kb._validate_import_url("http://127.0.0.1/x") is False
        assert kb._validate_import_url("http://192.168.1.1/x") is False
        assert kb._validate_import_url("http://10.0.0.5/x") is False
        assert kb._validate_import_url("http://172.16.0.1/x") is False

    @pytest.mark.asyncio
    async def test_import_url_fetches_and_creates(self, monkeypatch, tmp_path):
        import socket

        from neurova.knowledge.repository import KnowledgeRepository

        repo = KnowledgeRepository(str(tmp_path / "kb"))
        repo._items.clear()
        monkeypatch.setattr(kb, "_get_repository", lambda agent_id="default": repo)
        monkeypatch.setattr(
            socket,
            "getaddrinfo",
            lambda host, port: [(socket.AF_INET, socket.SOCK_STREAM, 6, "", ("93.184.216.34", 443))],
        )

        html = b"<html><body><h1>Remote Title</h1><p>Remote body content</p></body></html>"
        monkeypatch.setattr(kb, "_fetch_url", lambda url: html)

        items, status = kb._import_file_data(html, "example.com/article.html", "default", {"user_id": "u1"})
        assert status == "html"
        assert len(items) == 1
        assert "Remote body content" in items[0]["content"]
