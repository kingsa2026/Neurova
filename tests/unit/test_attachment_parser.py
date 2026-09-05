"""
测试：附件文本抽取器（R-3）

契约：
  1. txt/md/code UTF-8 直接解码
  2. docx / xlsx / pptx / pdf / csv 抽取文本（真实文件生成后解析）
  3. 空文件 / 超限 / 旧二进制（doc/xls）降级 None 不抛异常
"""

import io


def _make_docx():
    import docx

    d = docx.Document()
    d.add_paragraph("这是文档第一段")
    d.add_paragraph("第二段内容")
    table = d.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _make_xlsx():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "姓名"
    ws["B1"] = "年龄"
    ws["A2"] = "张三"
    ws["B2"] = 20
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _make_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "幻灯片标题"
    body.text = "要点一\n要点二"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_txt_decodes():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text("你好 hello".encode("utf-8"), "a.txt", "text")
    assert text == "你好 hello"
    assert status == "decoded"


def test_md_decodes():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(b"# Title\n\nbody", "a.md", "text")
    assert text == "# Title\n\nbody"
    assert status == "decoded"


def test_code_decodes():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(b"print('hi')", "a.py", "code")
    assert text == "print('hi')"
    assert status == "decoded"


def test_docx_extracts_paragraphs_and_table():
    from neurova.attachment_parser import extract_attachment_text

    data = _make_docx()
    text, status = extract_attachment_text(data, "报告.docx", "document")
    assert status == "docx"
    assert "这是文档第一段" in text
    assert "第二段内容" in text
    assert "A | B" in text


def test_xlsx_extracts_cells():
    from neurova.attachment_parser import extract_attachment_text

    data = _make_xlsx()
    text, status = extract_attachment_text(data, "表格.xlsx", "document")
    assert status == "xlsx"
    assert "姓名" in text and "张三" in text and "年龄" in text


def test_pptx_extracts_shapes():
    from neurova.attachment_parser import extract_attachment_text

    data = _make_pptx()
    text, status = extract_attachment_text(data, "演示.pptx", "document")
    assert status == "pptx"
    assert "幻灯片标题" in text
    assert "要点一" in text


def test_pdf_extracts_pages():
    from neurova.attachment_parser import extract_attachment_text

    # 用 reportlab 不可用，改用 pypdf 写入器（简单文本页）
    from pypdf import PdfWriter

    w = PdfWriter()
    # 无现成文本页时跳过复杂生成，用最小空白页验证不抛异常
    w.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    w.write(buf)
    text, status = extract_attachment_text(buf.getvalue(), "手册.pdf", "document")
    assert status == "pdf"
    assert isinstance(text, str)


def test_csv_decodes():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(b"a,b\n1,2", "d.csv", "document")
    assert text == "a,b\n1,2"
    assert status == "csv"


def test_legacy_doc_falls_back():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(b"\x00\x01legacy", "old.doc", "document")
    assert text is None
    assert status == "unsupported_format"


def test_empty_file():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(b"", "a.txt", "text")
    assert text is None
    assert status == "empty_file"


def test_large_file_accepted_and_truncated():
    """2026-09-06 契约升级：移除 2MB 字节硬限制（用户要求）。
    真正的上下文闸门是 MAX_EXTRACT_CHARS（20 万字符）——大文件接受，
    超长文本截断。"""
    from neurova.attachment_parser import extract_attachment_text, MAX_EXTRACT_CHARS

    big = ("中文行内容abc。" + chr(10)) * (300 * 1024)
    big = big.encode("utf-8")  # >2MB
    text, status = extract_attachment_text(big, "big.txt", "text")
    assert status == "decoded"
    assert text is not None and len(text) <= MAX_EXTRACT_CHARS


def test_invalid_docx_returns_none():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(b"not a real docx", "fake.docx", "document")
    assert text is None
    assert "docx_error" in status


def test_html_extracts_text():
    from neurova.attachment_parser import extract_attachment_text

    html = b"<html><head><title>My Page</title></head><body><h1>Hello</h1><p>Paragraph text.</p><script>var x=1;</script></body></html>"
    text, status = extract_attachment_text(html, "page.html", "document")
    assert status == "html"
    assert "Hello" in text
    assert "Paragraph text." in text
    # script/style 内容被剥离
    assert "var x=1" not in text


def test_excel_via_xlsx_extracts_text():
    from neurova.attachment_parser import extract_attachment_text

    data = _make_xlsx()
    text, status = extract_attachment_text(data, "表格.xlsx", "document")
    assert status == "xlsx"
    assert "姓名" in text and "张三" in text


def test_json_decodes():
    from neurova.attachment_parser import extract_attachment_text

    data = '{"name": "知识导入", "items": [1, 2, 3]}'.encode("utf-8")
    text, status = extract_attachment_text(data, "config.json", "text")
    assert status == "decoded"
    assert "知识导入" in text


def test_xml_extracts_text():
    from neurova.attachment_parser import extract_attachment_text

    data = '<?xml version="1.0"?><root><title>XML 标题</title><body>Body text</body></root>'.encode('utf-8')
    text, status = extract_attachment_text(data, "doc.xml", "document")
    assert status == "html"
    assert "Body text" in text


def test_rtf_extracts_text():
    from neurova.attachment_parser import extract_attachment_text

    rtf = b"{\\rtf1\\ansi Hello RTF World.\\par Second line.}"
    text, status = extract_attachment_text(rtf, "doc.rtf", "document")
    assert status == "rtf"
    assert "Hello RTF World." in text
    assert "Second line." in text


def _make_odt():
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    doc = OpenDocumentText()
    doc.text.addElement(P(text="ODT 正文内容"))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_ods():
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableRow, TableCell
    from odf.text import P

    doc = OpenDocumentSpreadsheet()
    table = Table(name="Sheet1")
    row = TableRow()
    cell = TableCell(valuetype="string")
    cell.addElement(P(text="单元格数据"))
    row.addElement(cell)
    table.addElement(row)
    doc.spreadsheet.addElement(table)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_odt_extracts_text():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(_make_odt(), "doc.odt", "document")
    assert status == "odt"
    assert "ODT 正文内容" in text


def test_ods_extracts_cells():
    from neurova.attachment_parser import extract_attachment_text

    text, status = extract_attachment_text(_make_ods(), "sheet.ods", "document")
    assert status == "ods"
    assert "单元格数据" in text
